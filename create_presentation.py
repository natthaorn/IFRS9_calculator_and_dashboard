import os
import sys

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx not found. Installing...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def create_deck():
    prs = Presentation()
    
    # Set to 16:9 Widescreen standard
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    DARK_BG = RGBColor(12, 16, 27)      # #0C101B
    WHITE = RGBColor(243, 244, 246)     # #F3F4F6
    MUTED = RGBColor(156, 163, 175)     # #9CA3AF
    ACCENT_BLUE = RGBColor(59, 130, 246) # #3B82F6
    ACCENT_TEAL = RGBColor(16, 185, 129) # #10B981
    ACCENT_AMBER = RGBColor(245, 158, 11) # #F59E0B
    
    # Helper to apply dark background
    def set_dark_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    # Helper to add standard header
    def add_header(slide, title_text, category_text="IFRS 9 ECL PORTFOLIO"):
        # Category label
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category_text.upper()
        cat_p.font.name = "Outfit"
        cat_p.font.size = Pt(11)
        cat_p.font.bold = True
        cat_p.font.color.rgb = ACCENT_BLUE
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = "Outfit"
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = WHITE
        
        # Line separator decoration
        # We can add a thin horizontal line or colored accent block
        # For simplicity, we just use layout spacing

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide1)
    
    # Large Title text box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "IFRS 9 Expected Credit Loss (ECL)\nCalculator & Analytics Dashboard"
    p1.font.name = "Outfit"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.LEFT
    
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Bridging Credit Risk Quantitative Modeling & Modern Dashboard Reporting"
    p1_sub.font.name = "Outfit"
    p1_sub.font.size = Pt(18)
    p1_sub.font.color.rgb = ACCENT_BLUE
    p1_sub.space_before = Pt(15)
    
    # Footer info
    info_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.2))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "Portfolio Risk Analytics Portfolio Project\nTech Stack: Python (FastAPI, Pandas, SciPy, Scikit-Learn) | HTML5, Vanilla CSS, JS, Chart.js"
    p_info.font.name = "Outfit"
    p_info.font.size = Pt(13)
    p_info.font.color.rgb = MUTED
    p_info.space_before = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 2: Data Quality & Validation Gates
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide2)
    add_header(slide2, "Data Quality & Pre-Calculation Controls")
    
    body_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf2 = body_box.text_frame
    tf2.word_wrap = True
    
    def add_bullet(tf, bold_prefix, text, color=WHITE, pt_size=15):
        p = tf.add_paragraph()
        p.space_before = Pt(12)
        p.level = 0
        run_bold = p.add_run()
        run_bold.text = bold_prefix + " "
        run_bold.font.bold = True
        run_bold.font.size = Pt(pt_size)
        run_bold.font.name = "Outfit"
        run_bold.font.color.rgb = ACCENT_BLUE
        
        run_text = p.add_run()
        run_text.text = text
        run_text.font.size = Pt(pt_size)
        run_text.font.name = "Outfit"
        run_text.font.color.rgb = color
        
    add_bullet(tf2, "[Completeness Controls]", "Scans and flags missing credit ratings or exposure amounts on load, preventing execution failures.", WHITE, 16)
    add_bullet(tf2, "[Logical Consistency Checks]", "Rejects records with negative Days Past Due (DPD), negative interest rates (EIR), or outstanding balances exceeding limits by more than 1.05x.", WHITE, 16)
    add_bullet(tf2, "[Outlier Screening]", "Flags abnormal Loan-To-Value ratios (LTV > 300%) or high interest rates (EIR > 35%) reflecting data input errors.", WHITE, 16)
    add_bullet(tf2, "[Audit-Ready Quarantining]", "Automatically separates clean data from dirty data (97.5% clean rate on mock portfolio), logging errors in a dedicated interactive table.", WHITE, 16)

    # -------------------------------------------------------------
    # SLIDE 3: IFRS 9 ECL Engine & Core Math
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide3)
    add_header(slide3, "IFRS 9 ECL Calculation Engine")
    
    body_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf3 = body_box.text_frame
    tf3.word_wrap = True
    
    add_bullet(tf3, "• Stage 1 (12-Month ECL):", "Performing assets. Receives 12m PD losses discounted at EIR.", WHITE, 16)
    add_bullet(tf3, "• Stage 2 (Lifetime ECL - SICR):", "Triggered if DPD >= 30, rating downgrades, Relative PD Ratio >= 3.0x, or Absolute PD Difference >= 2.0% (avoids false-positives on low-risk loans).", WHITE, 16)
    add_bullet(tf3, "• Stage 3 (Lifetime ECL - Default):", "Impaired assets (Rating 10 or DPD >= 90).", WHITE, 16)
    add_bullet(tf3, "• Vasicek TTC to PiT Transformation:", "Transforms Through-The-Cycle PDs to Point-in-Time curves based on macro systematic factor (Y) using Merton-Vasicek single-factor model.", WHITE, 16)
    add_bullet(tf3, "• Marginal PD & Collateralized LGD:", "Projects conditional and marginal PD curves for annual discounting. Computes LGD weighted by secured portions with haircut deductions (Real Estate 20%, Vehicle 35%).", WHITE, 16)

    # -------------------------------------------------------------
    # SLIDE 4: Budget Forecasting & Stress Testing
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide4)
    add_header(slide4, "ECL Forecasting & Scenario Analysis")
    
    body_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf4 = body_box.text_frame
    tf4.word_wrap = True
    
    add_bullet(tf4, "[Multi-Scenario Projections]", "Models ECL charges over a 5-year budget horizon (2026-2030) under Base Case, Optimistic Case, and Pessimistic (Stress Case) macroeconomic paths.", WHITE, 16)
    add_bullet(tf4, "[Portfolio Amortization Decay]", "Assumes a natural 15% annual portfolio decay to reflect loan paydowns and shrinking remaining terms over the forecast horizon.", WHITE, 16)
    add_bullet(tf4, "[Stress Test Sensitivities]", "Under severe economic stress (Pessimistic scenario), the model dynamically shifts ratings and projects spikes in Stage 2/3 and ECL charges.", WHITE, 16)
    add_bullet(tf4, "[Capital & Reserves Planning]", "Provides risk department outputs to align loan-loss provisioning and capital adequacy reserves for budget planning activities.", WHITE, 16)

    # -------------------------------------------------------------
    # SLIDE 5: Model Monitoring & Validation
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_bg(slide5)
    add_header(slide5, "Model Monitoring & Recalibration controls")
    
    body_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf5 = body_box.text_frame
    tf5.word_wrap = True
    
    add_bullet(tf5, "[Model Discrimination (Gini/AUC)]", "Calculates ROC curves and Gini coefficients (~0.82) to measure the rating model's ability to rank-order default risk.", WHITE, 16)
    add_bullet(tf5, "[Hardcoded Population Stability Index (PSI)]", "Implemented step-by-step in pure Python to measure rating distribution shift between origination (baseline) and today. Classifies shifts: Green (<0.10), Amber (0.10-0.25), Red (>=0.25).", WHITE, 16)
    add_bullet(tf5, "[Audit-Safe Backtesting]", "Compares average predicted 12-month PiT PD against actual observed default rates by rating grade to verify calibration correctness.", WHITE, 16)
    add_bullet(tf5, "[Model Drift Mitigation]", "Provides actionable triggers for risk managers (e.g. Amber/Red PSI status triggers model recalibration warning).", WHITE, 16)

    # Save presentation
    output_path = "/Users/fhyfhy/Desktop/ifrs9-ecl-analytics/IFRS9_ECL_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
